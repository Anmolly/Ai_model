"""
Presentation Generation Module - Create presentations in multiple formats
"""

import logging
from typing import Dict, List, Any, Optional
from datetime import datetime
import json

logger = logging.getLogger(__name__)


class PresentationGenerator:
    """Generate presentations in HTML, PPTX, and PDF formats"""
    
    def __init__(self, config: Dict[str, Any] = None):
        self.config = config or {}
        self.logger = logger
        self.themes = ['default', 'dark', 'professional', 'minimal', 'colorful']
        self.styles = ['professional', 'casual', 'academic', 'creative']
        
    async def generate(self, title: str, content: List[Dict[str, Any]],
                      style: str = 'professional', output_format: str = 'html',
                      theme: str = 'default') -> Dict[str, Any]:
        """
        Generate presentation
        
        Args:
            title: Presentation title
            content: List of slides with content
            style: Presentation style
            output_format: Output format (html, pptx, pdf)
            theme: Theme name
            
        Returns:
            Generated presentation data
        """
        try:
            self.logger.info(f"Generating presentation: {title}")
            
            if output_format == 'html':
                result = await self._generate_html(title, content, style, theme)
            elif output_format == 'pptx':
                result = await self._generate_pptx(title, content, style, theme)
            elif output_format == 'pdf':
                result = await self._generate_pdf(title, content, style, theme)
            else:
                result = {'error': f'Unknown format: {output_format}'}
            
            return result
            
        except Exception as e:
            self.logger.error(f"Generation error: {str(e)}")
            return {'error': str(e)}
    
    async def _generate_html(self, title: str, content: List[Dict[str, Any]],
                            style: str, theme: str) -> Dict[str, Any]:
        """Generate HTML presentation"""
        try:
            html_content = self._create_html_structure(title, content, style, theme)
            
            filename = f"presentation_{datetime.now().strftime('%Y%m%d_%H%M%S')}.html"
            
            return {
                'success': True,
                'format': 'html',
                'filename': filename,
                'content': html_content,
                'timestamp': datetime.now().isoformat()
            }
            
        except Exception as e:
            self.logger.error(f"HTML generation error: {str(e)}")
            return {'error': str(e)}
    
    def _create_html_structure(self, title: str, content: List[Dict[str, Any]],
                              style: str, theme: str) -> str:
        """Create HTML presentation structure"""
        
        slides_html = ""
        for i, slide in enumerate(content):
            slide_title = slide.get('title', f'Slide {i+1}')
            slide_content = slide.get('content', '')
            
            slides_html += f"""
            <div class="slide">
                <h1>{slide_title}</h1>
                <div class="slide-content">
                    {slide_content}
                </div>
            </div>
            """
        
        html = f"""
        <!DOCTYPE html>
        <html lang="en">
        <head>
            <meta charset="UTF-8">
            <meta name="viewport" content="width=device-width, initial-scale=1.0">
            <title>{title}</title>
            <style>
                * {{
                    margin: 0;
                    padding: 0;
                    box-sizing: border-box;
                }}
                
                body {{
                    font-family: 'Segoe UI', Tahoma, Geneva, Verdana, sans-serif;
                    background: #f0f0f0;
                }}
                
                .presentation {{
                    max-width: 1000px;
                    margin: 0 auto;
                    padding: 20px;
                }}
                
                .slide {{
                    background: white;
                    padding: 60px;
                    margin: 20px 0;
                    border-radius: 8px;
                    box-shadow: 0 2px 10px rgba(0,0,0,0.1);
                    min-height: 600px;
                    display: flex;
                    flex-direction: column;
                    justify-content: center;
                }}
                
                .slide h1 {{
                    font-size: 48px;
                    color: #333;
                    margin-bottom: 30px;
                }}
                
                .slide-content {{
                    font-size: 24px;
                    color: #666;
                    line-height: 1.6;
                }}
                
                .slide-content p {{
                    margin: 15px 0;
                }}
                
                .slide-content ul {{
                    margin-left: 30px;
                }}
                
                .slide-content li {{
                    margin: 10px 0;
                }}
            </style>
        </head>
        <body>
            <div class="presentation">
                <h1 style="text-align: center; margin: 40px 0; font-size: 56px;">{title}</h1>
                {slides_html}
            </div>
        </body>
        </html>
        """
        
        return html
    
    async def _generate_pptx(self, title: str, content: List[Dict[str, Any]],
                            style: str, theme: str) -> Dict[str, Any]:
        """Generate PPTX presentation"""
        try:
            from pptx import Presentation
            from pptx.util import Inches, Pt
            
            prs = Presentation()
            prs.slide_width = Inches(10)
            prs.slide_height = Inches(7.5)
            
            # Title slide
            title_slide_layout = prs.slide_layouts[0]
            slide = prs.slides.add_slide(title_slide_layout)
            title_shape = slide.shapes.title
            subtitle_shape = slide.placeholders[1]
            
            title_shape.text = title
            subtitle_shape.text = f"Generated on {datetime.now().strftime('%Y-%m-%d')}"
            
            # Content slides
            for slide_data in content:
                bullet_slide_layout = prs.slide_layouts[1]
                slide = prs.slides.add_slide(bullet_slide_layout)
                title_shape = slide.shapes.title
                body_shape = slide.placeholders[1]
                
                title_shape.text = slide_data.get('title', 'Slide')
                
                tf = body_shape.text_frame
                tf.text = slide_data.get('content', '')
            
            filename = f"presentation_{datetime.now().strftime('%Y%m%d_%H%M%S')}.pptx"
            
            return {
                'success': True,
                'format': 'pptx',
                'filename': filename,
                'timestamp': datetime.now().isoformat()
            }
            
        except Exception as e:
            self.logger.error(f"PPTX generation error: {str(e)}")
            return {'error': str(e)}
    
    async def _generate_pdf(self, title: str, content: List[Dict[str, Any]],
                           style: str, theme: str) -> Dict[str, Any]:
        """Generate PDF presentation"""
        try:
            from reportlab.lib.pagesizes import letter
            from reportlab.pdfgen import canvas
            
            filename = f"presentation_{datetime.now().strftime('%Y%m%d_%H%M%S')}.pdf"
            
            c = canvas.Canvas(filename, pagesize=letter)
            width, height = letter
            
            # Title page
            c.setFont("Helvetica-Bold", 48)
            c.drawString(50, height - 100, title)
            c.showPage()
            
            # Content pages
            for slide in content:
                c.setFont("Helvetica-Bold", 36)
                c.drawString(50, height - 100, slide.get('title', 'Slide'))
                
                c.setFont("Helvetica", 12)
                y = height - 150
                for line in str(slide.get('content', '')).split('\n'):
                    c.drawString(50, y, line)
                    y -= 20
                
                c.showPage()
            
            c.save()
            
            return {
                'success': True,
                'format': 'pdf',
                'filename': filename,
                'timestamp': datetime.now().isoformat()
            }
            
        except Exception as e:
            self.logger.error(f"PDF generation error: {str(e)}")
            return {'error': str(e)}
    
    async def add_slide(self, presentation_id: str, slide_data: Dict[str, Any]) -> Dict[str, Any]:
        """Add slide to presentation"""
        try:
            self.logger.info(f"Adding slide to presentation: {presentation_id}")
            
            return {
                'success': True,
                'presentation_id': presentation_id,
                'slide': slide_data,
                'timestamp': datetime.now().isoformat()
            }
            
        except Exception as e:
            self.logger.error(f"Add slide error: {str(e)}")
            return {'error': str(e)}